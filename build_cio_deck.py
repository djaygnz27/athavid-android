from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import copy

# Colors
PSEG_BLUE = RGBColor(0x00, 0x33, 0x87)
PSEG_ORANGE = RGBColor(0xF7, 0x94, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
GREEN = RGBColor(0x00, 0x80, 0x00)
YELLOW = RGBColor(0xFF, 0xC0, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)

def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_text_to_shape(shape, text, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def slide_header(slide, title, subtitle=None):
    # Top bar
    bar = add_rect(slide, 0, 0, 13.333, 1.0, fill_color=PSEG_BLUE)
    # Orange accent
    add_rect(slide, 0, 0.85, 13.333, 0.12, fill_color=PSEG_ORANGE)
    # Title text
    add_textbox(slide, title, 0.3, 0.12, 10, 0.65, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.6, 10, 0.35, font_size=12, color=PSEG_ORANGE, align=PP_ALIGN.LEFT)
    # Bottom bar
    add_rect(slide, 0, 7.2, 13.333, 0.3, fill_color=PSEG_BLUE)
    add_textbox(slide, "PSEG Long Island | JMUX Replacement Program | Confidential | March 2026", 
                0.3, 7.22, 12, 0.25, font_size=8, color=WHITE)

prs = new_prs()
blank_layout = prs.slide_layouts[6]  # Blank

# ─────────────────────────────────────────────
# SLIDE 1: TITLE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=PSEG_BLUE)
add_rect(slide, 0, 5.5, 13.333, 0.18, fill_color=PSEG_ORANGE)
add_textbox(slide, "JMUX Replacement Program", 1.0, 1.5, 11, 1.0, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "CIO Briefing | Phase 1 Complete · Phase 2 Underway · Path to Completion", 1.0, 2.6, 11, 0.6, font_size=16, color=PSEG_ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, "PRJ13797  |  Program Manager: Dhananjaya Gunaratne  |  March 12, 2026", 1.0, 3.4, 11, 0.5, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "PSEG Long Island  |  IT Infrastructure", 1.0, 4.2, 11, 0.5, font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(slide, "CONFIDENTIAL", 1.0, 6.8, 11, 0.4, font_size=10, color=PSEG_ORANGE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2: EXECUTIVE SUMMARY
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Executive Summary", "JMUX Replacement Program — Status as of March 12, 2026")
add_rect(slide, 0.3, 1.1, 12.7, 5.9, fill_color=LIGHT_GRAY)

# Three columns
cols = [
    ("✅  Phase 1\nCOMPLETE", "31 substations cutover to MPLS\nAll Wave 1 & Wave 2 cutovers done\nNSP installed at Melville & Hicksville\nSSP & CMDB documentation approved\nMinor close-outs in progress (Hicksville cabling, OTDR)", GREEN),
    ("🟡  Phase 2\nIN PROGRESS", "36-site engineering design underway\nSite surveys completed during Phase 1\nBurns & McDonnell issuing IFRs & IFCs\nNokia Phase 2 equipment delivered 2/19\nSchedule: Yellow (2.5-wk delay, recovering)", YELLOW),
    ("📋  CIO Actions\nREQUIRED", "Awareness: Schedule is yellow, recovery plan in place\nApproval: Confirm Phase 2 vendor engagement continues\nDecision: Resource support for SSP/security reviews\nSupport: Crown Castle fiber issue at Hicksville", PSEG_BLUE),
]
for i, (title, body, color) in enumerate(cols):
    x = 0.4 + i * 4.2
    add_rect(slide, x, 1.2, 4.0, 0.55, fill_color=color)
    add_textbox(slide, title, x+0.1, 1.22, 3.8, 0.5, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 1.75, 4.0, 5.1, fill_color=WHITE, line_color=color, line_width=1.5)
    add_textbox(slide, body, x+0.15, 1.85, 3.7, 4.8, font_size=11, color=DARK_GRAY)

# ─────────────────────────────────────────────
# SLIDE 3: PHASE 1 ACCOMPLISHMENTS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Phase 1 Accomplishments", "31 Substations | Complete")

add_rect(slide, 0.3, 1.1, 8.2, 5.9, fill_color=LIGHT_GRAY)
add_rect(slide, 8.7, 1.1, 4.3, 5.9, fill_color=LIGHT_GRAY)

# Left: accomplishments
add_rect(slide, 0.3, 1.1, 8.2, 0.45, fill_color=PSEG_BLUE)
add_textbox(slide, "KEY ACCOMPLISHMENTS", 0.4, 1.14, 8.0, 0.38, font_size=12, bold=True, color=WHITE)

items = [
    ("All 31 substation sites installed & commissioned", GREEN),
    ("Wave 1 cutover (12/8–12/18/25) — non-tie sites complete", GREEN),
    ("Wave 2 cutover (1/6–1/8/26) — tie sites complete", GREEN),
    ("CWDM cutover complete at Hicksville & ACC", GREEN),
    ("Nokia NSP racked & powered at Melville & Hicksville", GREEN),
    ("Staging completed at Hawkeye Warehouse — 9/29/25", GREEN),
    ("SSP documentation approved by security team", GREEN),
    ("CMDB inventory updated & confirmed — 2/24/26", GREEN),
    ("Phase 2 site surveys completed during Phase 1 (ahead of schedule)", GREEN),
    ("7 additional nodes added to complete Ring 1 & Ring 2", GREEN),
]
for j, (txt, col) in enumerate(items):
    y = 1.65 + j * 0.48
    add_rect(slide, 0.35, y, 0.25, 0.3, fill_color=col)
    add_textbox(slide, txt, 0.7, y-0.02, 7.6, 0.38, font_size=11, color=DARK_GRAY)

# Right: stats box
add_rect(slide, 8.7, 1.1, 4.3, 0.45, fill_color=PSEG_ORANGE)
add_textbox(slide, "PHASE 1 AT A GLANCE", 8.8, 1.14, 4.1, 0.38, font_size=12, bold=True, color=WHITE)
stats = [
    ("31 / 31", "Sites Complete"),
    ("100%", "Cutover Done"),
    ("2", "Network Rings Closed"),
    ("Sept 2025", "Staging Complete"),
    ("Jan 2026", "Final Cutover"),
    ("$4.44M", "Spent to Date"),
]
for j, (num, label) in enumerate(stats):
    y = 1.65 + j * 0.85
    add_rect(slide, 8.75, y, 4.2, 0.75, fill_color=WHITE, line_color=PSEG_ORANGE, line_width=1)
    add_textbox(slide, num, 8.8, y+0.02, 4.1, 0.38, font_size=18, bold=True, color=PSEG_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, 8.8, y+0.38, 4.1, 0.3, font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4: PHASE 2 STATUS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Phase 2 Status", "36 Sites | Engineering Design Underway | Burns & McDonnell SOW: $1.37M")

# Status bar
statuses = [
    ("Site Surveys", "COMPLETE", GREEN),
    ("Nokia Equip. Delivery", "COMPLETE", GREEN),
    ("IFR Design (B&McD)", "IN PROGRESS", YELLOW),
    ("IFC Design", "IN PROGRESS", YELLOW),
    ("Construction", "NOT STARTED", DARK_GRAY),
    ("Cutover", "NOT STARTED", DARK_GRAY),
]
for i, (label, status, color) in enumerate(statuses):
    x = 0.3 + i * 2.15
    add_rect(slide, x, 1.15, 2.0, 0.55, fill_color=color)
    add_textbox(slide, label, x+0.05, 1.16, 1.9, 0.28, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, status, x+0.05, 1.42, 1.9, 0.25, font_size=8, color=WHITE, align=PP_ALIGN.CENTER)

# Milestone table
add_rect(slide, 0.3, 1.85, 12.7, 0.38, fill_color=PSEG_BLUE)
headers = ["Milestone", "Contract Date", "Status", "Notes"]
widths = [4.0, 2.2, 2.0, 4.5]
x_pos = 0.35
for h, w in zip(headers, widths):
    add_textbox(slide, h, x_pos, 1.87, w, 0.32, font_size=11, bold=True, color=WHITE)
    x_pos += w

milestones = [
    ("Phase 1 Staging Support", "8/22/2025", "✅ Complete", "$57,300 paid"),
    ("Phase 1 Construction Coordinator", "Aug 2025–Jan 2026", "✅ Complete", "$368,100 — monthly invoiced"),
    ("Substation Site Surveys (Phase 2)", "2/27/2026", "✅ Complete", "Completed during Phase 1 — ahead of schedule"),
    ("MPLS Network Node Design — IFR", "Wave 1: 2/23/26", "🟡 In Progress", "B&McD issuing IFRs; PSEG LI review 2-wk window"),
    ("MPLS Network Node Design — IFC", "Wave 1: 3/23/26", "⏳ Upcoming", "IFC follows PSEG LI review & comment"),
    ("Construction & Installation", "Aug–Dec 2026", "⏳ Not Started", "Pending IFC completion"),
]
for j, (m, date, status, notes) in enumerate(milestones):
    y = 2.3 + j * 0.72
    bg = LIGHT_GRAY if j % 2 == 0 else WHITE
    add_rect(slide, 0.3, y, 12.7, 0.68, fill_color=bg)
    row = [m, date, status, notes]
    x_pos = 0.35
    for val, w in zip(row, widths):
        col = GREEN if "Complete" in val else (YELLOW if "Progress" in val else DARK_GRAY)
        add_textbox(slide, val, x_pos, y+0.05, w-0.1, 0.58, font_size=10, color=col if "✅" in val or "🟡" in val else DARK_GRAY)
        x_pos += w

# ─────────────────────────────────────────────
# SLIDE 5: SCHEDULE STATUS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Schedule Status", "Overall: 🟡 Yellow | Recovery Plan in Place")

add_rect(slide, 0.3, 1.1, 5.8, 5.9, fill_color=LIGHT_GRAY)
add_rect(slide, 6.3, 1.1, 6.7, 5.9, fill_color=LIGHT_GRAY)

# Left: SPI gauge visual
add_rect(slide, 0.3, 1.1, 5.8, 0.45, fill_color=YELLOW)
add_textbox(slide, "SCHEDULE PERFORMANCE INDEX (SPI)", 0.4, 1.13, 5.6, 0.38, font_size=11, bold=True, color=DARK_GRAY)

add_textbox(slide, "SPI: 0.75", 0.5, 1.7, 5.5, 0.8, font_size=36, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_textbox(slide, "🟡  YELLOW — Schedule delayed ~2.5 weeks", 0.5, 2.55, 5.5, 0.5, font_size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_textbox(slide, "Root Causes:", 0.5, 3.15, 5.4, 0.35, font_size=12, bold=True, color=DARK_GRAY)
causes = [
    "SSP documentation delays (security review queue)",
    "CMDB update coordination time",
    "Inclement weather affecting field work",
    "Crown Castle fiber issue at Hicksville (OTDR testing pending)",
]
for j, c in enumerate(causes):
    add_textbox(slide, f"• {c}", 0.5, 3.55 + j*0.45, 5.4, 0.42, font_size=10, color=DARK_GRAY)

add_textbox(slide, "Recovery Actions:", 0.5, 5.55, 5.4, 0.35, font_size=12, bold=True, color=PSEG_BLUE)
add_textbox(slide, "All approvals received. Remaining Phase 1 close-out is remote configuration — no physical site access needed. Phase 2 IFR/IFC design on track.", 0.5, 5.95, 5.4, 0.9, font_size=10, color=DARK_GRAY)

# Right: timeline
add_rect(slide, 6.3, 1.1, 6.7, 0.45, fill_color=PSEG_BLUE)
add_textbox(slide, "PROGRAM TIMELINE", 6.4, 1.13, 6.5, 0.38, font_size=11, bold=True, color=WHITE)

timeline = [
    ("Apr 2024", "Project Start", GREEN, True),
    ("Oct 2024", "RFP Award & SOW Executed", GREEN, True),
    ("Jan 2025", "Phase 1 Detailed Design Start", GREEN, True),
    ("Sep 2025", "Staging Complete (Hawkeye)", GREEN, True),
    ("Oct 2025", "Phase 1 Construction Start", GREEN, True),
    ("Dec 2025 – Jan 2026", "Wave 1 & 2 Cutovers Complete", GREEN, True),
    ("Feb 2026", "Nokia Phase 2 Equipment Delivered", GREEN, True),
    ("Mar 2026", "Phase 2 IFR Design in Progress  ◀ NOW", PSEG_ORANGE, False),
    ("Mar–Jul 2026", "Phase 2 IFC Design & Review", DARK_GRAY, False),
    ("Aug–Dec 2026", "Phase 2 Construction", DARK_GRAY, False),
    ("Jul 2027", "Program Completion (Target)", DARK_GRAY, False),
]
for j, (date, event, color, done) in enumerate(timeline):
    y = 1.65 + j * 0.5
    add_rect(slide, 6.35, y, 0.2, 0.28, fill_color=color)
    add_textbox(slide, f"{date}  —  {event}", 6.65, y-0.02, 6.2, 0.38, font_size=10, bold=(not done), color=color if not done else DARK_GRAY)

# ─────────────────────────────────────────────
# SLIDE 6: FINANCIALS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Financial Status", "Program Budget: $6,600,000 | Status: ON BUDGET ✅")

# Big numbers row
fin_boxes = [
    ("$6,600,000", "Total URB\nApproved Budget", PSEG_BLUE),
    ("$4,437,070", "Spent\nTo Date (67%)", GREEN),
    ("$2,162,930", "Estimate to\nComplete", PSEG_ORANGE),
    ("$0", "Over /\nUnder Budget", GREEN),
]
for i, (num, label, color) in enumerate(fin_boxes):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 1.15, 3.0, 1.6, fill_color=color)
    add_textbox(slide, num, x+0.1, 1.25, 2.8, 0.8, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x+0.1, 2.0, 2.8, 0.65, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Progress bar
add_textbox(slide, "Budget Utilization", 0.3, 2.95, 5, 0.35, font_size=12, bold=True, color=DARK_GRAY)
add_rect(slide, 0.3, 3.35, 12.7, 0.45, fill_color=RGBColor(0xDD, 0xDD, 0xDD))
add_rect(slide, 0.3, 3.35, 12.7 * 0.67, 0.45, fill_color=GREEN)
add_textbox(slide, "67% Spent", 0.5, 3.38, 8, 0.35, font_size=12, bold=True, color=WHITE)
add_textbox(slide, "33% Remaining", 9.0, 3.38, 4, 0.35, font_size=12, bold=True, color=DARK_GRAY)

# PO breakdown
add_rect(slide, 0.3, 3.95, 12.7, 0.38, fill_color=PSEG_BLUE)
po_headers = ["PO #", "Supplier", "PO Value", "Invoiced", "Remaining", "Status"]
po_widths = [1.5, 3.5, 1.8, 1.8, 1.8, 2.3]
x_pos = 0.35
for h, w in zip(po_headers, po_widths):
    add_textbox(slide, h, x_pos, 3.97, w, 0.32, font_size=10, bold=True, color=WHITE)
    x_pos += w

po_data = [
    ("7000001429", "Burns & McDonnell", "$137,500", "$49,500", "$88,000", "🟡 Active"),
    ("7000001415", "Burns & McDonnell", "$869,235", "$798,665", "$70,570", "🟡 Active"),
    ("7000002099", "Burns & McDonnell (Ph2)", "$1,370,040", "$388,740", "$981,300", "🟡 Active"),
    ("7000001857", "Nokia of America", "$1,899,030", "$532,525", "$1,366,505", "🟡 Active"),
    ("7000002054", "Elecnor Hawkeye", "$402,793", "$257,619", "$145,174", "🟡 Active"),
    ("7000001175", "Nokia of America", "$190,575", "$190,575", "$0", "✅ Paid in Full"),
    ("7000001748", "Presidio Networked Sol.", "$40,955", "$40,955", "$0", "✅ Paid in Full"),
    ("7000001767", "Graybar Electric", "$24,759", "$24,759", "$0", "✅ Paid in Full"),
]
for j, row in enumerate(po_data):
    y = 4.4 + j * 0.37
    bg = LIGHT_GRAY if j % 2 == 0 else WHITE
    add_rect(slide, 0.3, y, 12.7, 0.35, fill_color=bg)
    x_pos = 0.35
    for val, w in zip(row, po_widths):
        col = GREEN if "Paid" in val else (YELLOW if "Active" in val else DARK_GRAY)
        add_textbox(slide, val, x_pos, y+0.02, w-0.1, 0.3, font_size=9, color=DARK_GRAY)
        x_pos += w

# ─────────────────────────────────────────────
# SLIDE 7: RISKS & ISSUES
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Risks & Issues", "Active Risks | Mitigations in Place")

add_rect(slide, 0.3, 1.1, 6.0, 0.4, fill_color=RED)
add_textbox(slide, "⚠  ACTIVE RISKS", 0.4, 1.13, 5.8, 0.32, font_size=12, bold=True, color=WHITE)
add_rect(slide, 6.5, 1.1, 6.5, 0.4, fill_color=PSEG_ORANGE)
add_textbox(slide, "🔧  ACTIVE ISSUES / IN RESOLUTION", 6.6, 1.13, 6.3, 0.32, font_size=12, bold=True, color=WHITE)

risks = [
    ("Legacy Equipment Fragility", "HIGH", "JMUX equipment 25+ yrs old. Faulty connectors, brittle fiber panels. Potential for unplanned outage during cutover.", "Careful cutover planning; dedicated B&McD on-site support. Phase 1 complete — risk reduced going into Ph2."),
    ("Substation Access Delays", "MEDIUM", "Weather, scheduling conflicts, or escort availability can delay field work.", "Verify weekly before mobilizing. Adjust site plan if specific sites inaccessible."),
    ("Phase 2 Schedule Slippage", "MEDIUM", "IFR/IFC design cycle (4-wk IFR + 2-wk review + 2-wk IFC) leaves limited buffer.", "Phase 2 surveys done early. B&McD running concurrent site designs. Weekly status calls tracking."),
    ("Crown Castle Fiber Issue (Hicksville)", "LOW", "Low dB levels on fiber — may be Crown Castle's network, not PSEG LI.", "OTDR testing planned 3/3. If Crown Castle issue confirmed, escalate to vendor."),
]
for j, (title, sev, desc, mit) in enumerate(risks):
    y = 1.6 + j * 1.35
    sev_col = RED if sev == "HIGH" else (YELLOW if sev == "MEDIUM" else GREEN)
    add_rect(slide, 0.3, y, 6.0, 1.25, fill_color=WHITE, line_color=sev_col, line_width=2)
    add_rect(slide, 0.3, y, 1.1, 0.38, fill_color=sev_col)
    add_textbox(slide, sev, 0.32, y+0.04, 1.0, 0.3, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 1.45, y+0.04, 4.7, 0.32, font_size=11, bold=True, color=DARK_GRAY)
    add_textbox(slide, desc, 0.4, y+0.42, 5.7, 0.42, font_size=9, color=DARK_GRAY)
    add_textbox(slide, f"Mitigation: {mit}", 0.4, y+0.82, 5.7, 0.38, font_size=9, italic=True, color=PSEG_BLUE)

issues = [
    ("Cat 6 Cabling — Hicksville NSP", "IN PROGRESS", "Cables ordered; Hawkeye scheduled to complete 3/3/2026."),
    ("OTDR Testing — Hicksville", "SCHEDULED", "Low dB levels. OTDR test set confirmed for 1550 & 1310nm. Site visit 3/3–3/4/2026."),
    ("Remote Access — Phase 2 Nodes", "IN PROGRESS", "NSP at Hicksville & Melville. Remote access needed before router configuration can begin."),
    ("Phase 2 SSP / Security Reviews", "COMPLETE", "Architecture review, security assessment, compliance readiness all approved 2/10/2026."),
]
for j, (title, status, desc) in enumerate(issues):
    y = 1.6 + j * 1.35
    st_col = GREEN if status == "COMPLETE" else (YELLOW if status == "IN PROGRESS" else PSEG_BLUE)
    add_rect(slide, 6.5, y, 6.5, 1.25, fill_color=WHITE, line_color=st_col, line_width=2)
    add_rect(slide, 6.5, y, 1.8, 0.38, fill_color=st_col)
    add_textbox(slide, status, 6.52, y+0.04, 1.7, 0.3, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 8.4, y+0.04, 4.4, 0.32, font_size=11, bold=True, color=DARK_GRAY)
    add_textbox(slide, desc, 6.6, y+0.48, 6.2, 0.7, font_size=10, color=DARK_GRAY)

# ─────────────────────────────────────────────
# SLIDE 8: CIO DECISIONS & ASKS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "CIO Actions Required", "3 Items for Decision or Awareness")

asks = [
    (
        "1",
        "AWARENESS",
        "Schedule Recovery — Yellow Status",
        YELLOW,
        "Phase 1 experienced a 2.5-week delay due to SSP documentation processing and inclement weather. The SPI stands at 0.75 for February.",
        "Recovery plan is in place. Remaining Phase 1 close-out items are remote (no site access required). Phase 2 engineering is running concurrently. Overall program target end date remains July 2027.",
        "No action required — for awareness only."
    ),
    (
        "2",
        "APPROVAL",
        "Continue Phase 2 Vendor Engagement",
        PSEG_BLUE,
        "Burns & McDonnell Phase 2 SOW ($1.37M) is active. IFR design has commenced for 36 sites. Nokia equipment delivered to staging. Phase 2 is on track to hit IFC milestones per SOW schedule.",
        "Confirm continued engagement of B&McD and Elecnor Hawkeye for Phase 2 construction support as design packages are issued.",
        "Approval requested to proceed with Phase 2 construction mobilization once IFC packages are complete (expected Q2 2026)."
    ),
    (
        "3",
        "SUPPORT",
        "Crown Castle Fiber Issue — Hicksville",
        RED,
        "OTDR testing revealed low dB levels on fiber at Hicksville substation. This may be a Crown Castle infrastructure issue outside PSEG LI's control.",
        "OTDR testing was scheduled for 3/3–3/4/2026. If confirmed as Crown Castle issue, formal escalation to Crown Castle will be required, which may impact Phase 2 schedule at that site.",
        "Request CIO support if vendor escalation to Crown Castle is needed to expedite resolution."
    ),
]

for j, (num, badge, title, color, context, detail, ask) in enumerate(asks):
    y = 1.15 + j * 1.95
    add_rect(slide, 0.3, y, 12.7, 1.85, fill_color=WHITE, line_color=color, line_width=2)
    add_rect(slide, 0.3, y, 0.65, 1.85, fill_color=color)
    add_textbox(slide, num, 0.3, y+0.65, 0.65, 0.55, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.05, y+0.05, 1.4, 0.38, fill_color=color)
    add_textbox(slide, badge, 1.07, y+0.08, 1.35, 0.3, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 2.55, y+0.07, 9.5, 0.38, font_size=14, bold=True, color=DARK_GRAY)
    add_textbox(slide, f"Context: {context}", 1.05, y+0.52, 11.8, 0.42, font_size=10, color=DARK_GRAY)
    add_textbox(slide, f"Detail: {detail}", 1.05, y+0.95, 11.8, 0.42, font_size=10, color=DARK_GRAY)
    add_textbox(slide, f"➤  {ask}", 1.05, y+1.42, 11.8, 0.35, font_size=10, bold=True, color=color)

# ─────────────────────────────────────────────
# SLIDE 9: NEXT STEPS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Next Steps & Path to Completion", "March 2026 — July 2027")

add_rect(slide, 0.3, 1.15, 5.9, 5.9, fill_color=LIGHT_GRAY)
add_rect(slide, 6.4, 1.15, 6.6, 5.9, fill_color=LIGHT_GRAY)

add_rect(slide, 0.3, 1.15, 5.9, 0.42, fill_color=PSEG_BLUE)
add_textbox(slide, "IMMEDIATE (Next 30 Days)", 0.4, 1.18, 5.7, 0.35, font_size=12, bold=True, color=WHITE)

near = [
    ("Mar 3", "Hawkeye completes Cat 6 cabling for NSP at Hicksville"),
    ("Mar 3–4", "OTDR testing at Hicksville — resolve Crown Castle fiber issue"),
    ("Mar 2026", "Remote router configuration begins (Layer 1-3 setup)"),
    ("Mar 2026", "B&McD continues Wave 1 IFR submissions (Far Rockaway, Valley Stream, Cedarhurst)"),
    ("Mar 9", "PSEG LI review window opens for Wave 1 IFR packages"),
    ("Mar 23", "Wave 1 IFC target completion"),
    ("Ongoing", "Weekly status calls with B&McD; monthly schedule updates"),
]
for j, (date, task) in enumerate(near):
    y = 1.68 + j * 0.62
    add_rect(slide, 0.35, y, 1.1, 0.48, fill_color=PSEG_ORANGE)
    add_textbox(slide, date, 0.37, y+0.06, 1.05, 0.35, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, task, 1.55, y+0.06, 4.5, 0.42, font_size=10, color=DARK_GRAY)

add_rect(slide, 6.4, 1.15, 6.6, 0.42, fill_color=PSEG_BLUE)
add_textbox(slide, "PATH TO COMPLETION", 6.5, 1.18, 6.4, 0.35, font_size=12, bold=True, color=WHITE)

path = [
    ("Q1 2026", "Phase 1 close-out & Phase 2 IFR design", PSEG_ORANGE),
    ("Q2 2026", "Phase 2 IFC completion (all 36 sites)", PSEG_ORANGE),
    ("Q2–Q3 2026", "Phase 2 construction mobilization", PSEG_BLUE),
    ("Q3–Q4 2026", "Phase 2 installations at all 36 sites", PSEG_BLUE),
    ("Q4 2026", "Phase 2 cutovers — Waves 3 & 4", PSEG_BLUE),
    ("Q1 2027", "As-built documentation complete", DARK_GRAY),
    ("Jul 2027", "🏁  Program Complete", GREEN),
]
for j, (period, task, color) in enumerate(path):
    y = 1.68 + j * 0.72
    add_rect(slide, 6.45, y, 6.45, 0.62, fill_color=WHITE, line_color=color, line_width=1.5)
    add_rect(slide, 6.45, y, 1.5, 0.62, fill_color=color)
    add_textbox(slide, period, 6.47, y+0.12, 1.45, 0.38, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, task, 8.05, y+0.12, 4.75, 0.38, font_size=11, color=DARK_GRAY)

# Save
prs.save('/app/cio_jmux_presentation.pptx')
print("✅ Presentation saved successfully!")
EOF
