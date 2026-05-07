from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ──────────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x0C, 0x1A)
GOLD      = RGBColor(0xF5, 0xC8, 0x42)
CORAL     = RGBColor(0xFF, 0x6B, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARKGOLD  = RGBColor(0xC8, 0x9F, 0x20)
SOFTWHITE = RGBColor(0xF0, 0xF0, 0xF0)
MIDNAVY   = RGBColor(0x16, 0x18, 0x30)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=NAVY):
    left = top = Inches(0)
    w = prs.slide_width
    h = prs.slide_height
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_bar(slide, color=GOLD, height=Inches(0.07), top=Inches(6.9)):
    shape = slide.shapes.add_shape(1, Inches(0), top, prs.slide_width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def left_bar(slide, color=GOLD, width=Inches(0.12), top=Inches(1.4), height=Inches(4.5)):
    shape = slide.shapes.add_shape(1, Inches(0.5), top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb

def divider(slide, top, color=GOLD, opacity=None):
    shape = slide.shapes.add_shape(1, Inches(0.7), top, Inches(11.9), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def bullet_block(slide, items, left, top, width, font_size=16, color=WHITE, spacing=0.42):
    for i, item in enumerate(items):
        add_text(slide, item, left, top + Inches(i * spacing), width, Inches(0.4),
                 font_size=font_size, color=color)

def stat_box(slide, value, label, left, top, w=Inches(2.3), h=Inches(1.4), val_color=GOLD):
    box = slide.shapes.add_shape(1, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = GOLD
    box.line.width = Pt(1.2)
    add_text(slide, value, left, top + Inches(0.08), w, Inches(0.65),
             font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.72), w, Inches(0.55),
             font_size=12, color=LIGHTGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
# Gold diagonal accent
shape = s.shapes.add_shape(1, Inches(9.5), Inches(0), Inches(3.83), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x10, 0x11, 0x22)
shape.line.fill.background()

accent_bar(s, GOLD, height=Inches(0.05), top=Inches(0))
accent_bar(s, GOLD, height=Inches(0.05), top=Inches(7.45))

add_text(s, "SACHI STREAM", Inches(0.8), Inches(1.6), Inches(8.5), Inches(1.1),
         font_size=54, bold=True, color=GOLD)
add_text(s, "Creator Economy Platform", Inches(0.8), Inches(2.75), Inches(8.5), Inches(0.6),
         font_size=26, bold=False, color=WHITE)
add_text(s, "Series A Investment Opportunity", Inches(0.8), Inches(3.45), Inches(8.5), Inches(0.5),
         font_size=18, color=LIGHTGRAY, italic=True)

divider(s, Inches(4.1), GOLD)

add_text(s, "$5,000,000 Seed Round", Inches(0.8), Inches(4.3), Inches(8), Inches(0.55),
         font_size=22, bold=True, color=CORAL)
add_text(s, "LDNA Consulting LLC  ·  sachistream.com  ·  May 2026",
         Inches(0.8), Inches(5.05), Inches(9), Inches(0.4),
         font_size=13, color=LIGHTGRAY)

add_text(s, "Confidential\nNot for Distribution", Inches(10.0), Inches(6.5), Inches(3), Inches(0.7),
         font_size=10, color=LIGHTGRAY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE VISION
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "The Vision", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "A global short-video and live creator platform built for the next generation of digital storytellers — with monetization, community, and cultural authenticity at its core.",
         Inches(0.85), Inches(1.25), Inches(11.5), Inches(1.0),
         font_size=20, color=WHITE)

add_text(s, "What We're Building", Inches(0.85), Inches(2.5), Inches(10), Inches(0.45),
         font_size=17, bold=True, color=GOLD)

bullets = [
    "🎬  Short-form video, photo carousels, and text posts — multi-format creator toolkit",
    "🎙️  Live podcasting and creator rooms with real-time gifting and audience engagement",
    "💰  In-app coin economy — creators earn directly from fans, no algorithm dependency",
    "🌍  Location-aware content, cultural hashtag discovery, and global reach from day one",
    "🛡️  AI-powered content moderation with mature content gating and NERC-compliant infrastructure",
]
bullet_block(s, bullets, Inches(0.85), Inches(3.05), Inches(11.8), font_size=15, spacing=0.44)

add_text(s, "Sachi (幸) — Japanese for happiness, good fortune, and bliss.",
         Inches(0.85), Inches(6.1), Inches(11), Inches(0.45),
         font_size=13, italic=True, color=LIGHTGRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM & OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, CORAL)

add_text(s, "The Problem & Opportunity", Inches(0.85), Inches(0.3), Inches(11), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1), CORAL)

# Two columns
add_text(s, "THE PROBLEM", Inches(0.85), Inches(1.3), Inches(5.5), Inches(0.4),
         font_size=16, bold=True, color=CORAL)
probs = [
    "✗  TikTok faces US ban risk — 170M US users need a home",
    "✗  Creators earn pennies from ad revenue splits",
    "✗  Algorithm-driven platforms suppress organic reach",
    "✗  No platform combines video + podcasting + live gifting",
    "✗  Cultural and diaspora communities underserved",
]
bullet_block(s, probs, Inches(0.85), Inches(1.8), Inches(5.8), font_size=14, spacing=0.44)

add_text(s, "THE OPPORTUNITY", Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.4),
         font_size=16, bold=True, color=GOLD)
opps = [
    "✓  $300B+ global creator economy by 2030",
    "✓  1B+ short-video viewers globally, growing 20% YoY",
    "✓  Live streaming gifts market: $20B+ and accelerating",
    "✓  Diaspora communities — 280M people, underserved",
    "✓  First-mover: no platform combines all these features",
]
bullet_block(s, opps, Inches(7.2), Inches(1.8), Inches(5.8), font_size=14, spacing=0.44)

# Vertical divider
shape = s.shapes.add_shape(1, Inches(6.85), Inches(1.3), Inches(0.03), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHTGRAY
shape.line.fill.background()

add_text(s, "\"The window to build the next dominant social platform is open right now. Sachi Stream is ready.\"",
         Inches(0.85), Inches(6.05), Inches(11.5), Inches(0.55),
         font_size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT WE'VE BUILT (PRODUCT)
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "What We've Built", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "Fully functional platform — live at sachistream.com — built and operational as of May 2026",
         Inches(0.85), Inches(1.2), Inches(11.5), Inches(0.45),
         font_size=16, color=LIGHTGRAY, italic=True)

features = [
    ("🎬 Multi-Format Feed",        "Short video, photo carousels, text posts\nAI moderation + mature content gating"),
    ("🎙️ Live Podcasting",          "Creator rooms, RTMP streaming via\nCloudflare, OBS-ready, episode archive"),
    ("🔴 Live Rooms & Gifting",     "Real-time comments, guest requests,\nWebRTC co-hosting, coin-powered gifts"),
    ("💬 Social Graph",             "Follow, DM, notifications, bookmarks,\nblocking, threaded comments"),
    ("🪙 Coin Economy",             "5 coin packages, Stripe payments,\nhost earns 80% of gift face value"),
    ("🛡️ Moderation & Safety",      "AI detection, report queue, MOD panel,\nuser block, mature content toggles"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.7 + col * 4.15)
    top  = Inches(1.85 + row * 1.85)
    box = s.shapes.add_shape(1, left, top, Inches(3.85), Inches(1.65))
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = GOLD
    box.line.width = Pt(0.8)
    add_text(s, title, left + Inches(0.15), top + Inches(0.1), Inches(3.55), Inches(0.38),
             font_size=14, bold=True, color=GOLD)
    add_text(s, desc,  left + Inches(0.15), top + Inches(0.52), Inches(3.55), Inches(1.0),
             font_size=12, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — TRACTION & MILESTONES
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Traction & Milestones", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

# Stats row
stat_box(s, "LIVE",    "Platform Status\nMay 2026",        Inches(0.7),  Inches(1.3), val_color=CORAL)
stat_box(s, "27",      "Backend Functions\nDeployed",      Inches(3.2),  Inches(1.3))
stat_box(s, "50+",     "Founding Creator\nSpots Open",     Inches(5.7),  Inches(1.3))
stat_box(s, "Android", "Closed Beta\nGoogle Play",         Inches(8.2),  Inches(1.3))
stat_box(s, "™",       "Sachi Trademark\nUSPTO Filed",     Inches(10.7), Inches(1.3))

add_text(s, "Key Milestones Achieved", Inches(0.85), Inches(2.95), Inches(10), Inches(0.4),
         font_size=17, bold=True, color=GOLD)

milestones = [
    "Q4 2025  →  Platform concept, brand identity, and technical architecture finalized",
    "Q1 2026  →  Full-stack MVP built: video, live, podcasting, coins, DMs, notifications",
    "Apr 2026  →  Android APK submitted to Google Play closed testing track",
    "Apr 2026  →  Sachi™ word mark filed with USPTO (Class 41 — Entertainment Services)",
    "Apr 2026  →  Founding Creator program launched — 50 curated beta creators",
    "May 2026  →  Platform live at sachistream.com — open beta imminent",
]
bullet_block(s, milestones, Inches(0.85), Inches(3.45), Inches(11.8), font_size=14, spacing=0.47)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Business Model", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

streams = [
    ("🪙  Virtual Gifting & Coins",
     "Users purchase coin bundles ($1.99–$99.99)\nHosts receive 80% of gift value\nPlatform retains 20% — primary near-term revenue",
     GOLD),
    ("📺  Creator Subscriptions",
     "Monthly fan subscriptions to creator channels\nRevenue split: 85% creator / 15% Sachi\nLaunching Q3 2026",
     CORAL),
    ("🎯  Native Advertising",
     "Brand integrations and sponsored content\nTargeted by interest, location, and hashtag\nLaunching at 100K MAU threshold",
     RGBColor(0x7B, 0xD4, 0xFF)),
    ("🎙️  Podcast Premium",
     "Premium episode access and listener subscriptions\nPodcast hosting-as-a-service for creators\nLaunching Q4 2026",
     RGBColor(0xA0, 0xFF, 0xB0)),
]

for i, (title, desc, color) in enumerate(streams):
    left = Inches(0.7 + (i % 2) * 6.1)
    top  = Inches(1.5 + (i // 2) * 2.2)
    box = s.shapes.add_shape(1, left, top, Inches(5.7), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_text(s, title, left + Inches(0.18), top + Inches(0.1), Inches(5.3), Inches(0.4),
             font_size=15, bold=True, color=color)
    add_text(s, desc,  left + Inches(0.18), top + Inches(0.55), Inches(5.3), Inches(1.3),
             font_size=13, color=SOFTWHITE)

add_text(s, "Unit Economics: At 1M MAU — projected $4-8M ARR from gifting alone at 2% conversion rate",
         Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.45),
         font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — GROWTH STRATEGY
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Growth Strategy", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

phases = [
    ("PHASE 1  ·  0–6 Months", "Months 1–6",
     ["50 Founding Creators — seed authentic content library",
      "Diaspora community targeting: South Asian, Sri Lankan, Pacific Islander",
      "App Store launch (iOS + Android) — Q3 2026",
      "Influencer partnerships and creator referral program",
      "Target: 10,000 MAU"], GOLD),
    ("PHASE 2  ·  6–18 Months", "Months 7–18",
     ["Brand ambassador program in 5 key markets",
      "Podcast network partnerships — syndication deals",
      "Targeted paid acquisition: TikTok refugees, diaspora communities",
      "API open for third-party integrations",
      "Target: 250,000 MAU"], CORAL),
    ("PHASE 3  ·  18–36 Months", "Months 19–36",
     ["Bare metal infrastructure deployment — Colombo, Sri Lanka HQ",
      "Regional expansion: South Asia, Pacific, Middle East",
      "Creator fund launch: $1M/year distributed to top creators",
      "Series B raise based on MAU and revenue traction",
      "Target: 1M–5M MAU"], RGBColor(0x7B, 0xD4, 0xFF)),
]

for i, (title, sub, items, color) in enumerate(phases):
    left = Inches(0.7 + i * 4.15)
    top  = Inches(1.3)
    box = s.shapes.add_shape(1, left, top, Inches(3.9), Inches(5.3))
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_text(s, title, left + Inches(0.15), top + Inches(0.12), Inches(3.6), Inches(0.55),
             font_size=13, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(s, "· " + item, left + Inches(0.15), top + Inches(0.75 + j * 0.82),
                 Inches(3.6), Inches(0.75), font_size=12, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — GLOBAL OPERATIONS: COLOMBO HQ
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Global Operations: Colombo, Sri Lanka", Inches(0.85), Inches(0.3), Inches(11.5), Inches(0.65),
         font_size=34, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "A dedicated engineering and operations hub in Colombo, Sri Lanka — combining world-class talent, cost efficiency, and 24/7 global coverage.",
         Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.65),
         font_size=16, color=WHITE)

# Left column — team
add_text(s, "ENGINEERING TEAM", Inches(0.85), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=15, bold=True, color=GOLD)
team = [
    "· 10–15 full-stack developers (Phase 1 hire)",
    "· Specialties: React, Node.js, mobile (iOS/Android)",
    "· Dedicated QA and DevOps engineers",
    "· Led by a senior engineering director",
    "· Average Sri Lankan developer cost: $15–25K/yr",
    "  vs. $120–180K US equivalent — 6–8x cost advantage",
]
bullet_block(s, team, Inches(0.85), Inches(2.6), Inches(5.5), font_size=14, spacing=0.44)

# Right column — responsibilities
add_text(s, "RESPONSIBILITIES", Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=15, bold=True, color=CORAL)
resp = [
    "· Platform cleanup, optimization, and bug resolution",
    "· Feature development and roadmap execution",
    "· Migration management for infrastructure scaling",
    "· 24/7 customer service and platform monitoring",
    "· App Store compliance and release management",
    "· Future: on-site bare metal infrastructure ops",
]
bullet_block(s, resp, Inches(7.0), Inches(2.6), Inches(5.8), font_size=14, spacing=0.44)

# Bottom infrastructure note
box = s.shapes.add_shape(1, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.25))
box.fill.solid()
box.fill.fore_color.rgb = MIDNAVY
box.line.color.rgb = GOLD
box.line.width = Pt(1)
add_text(s, "🏢  Phase 3 Infrastructure: Bare metal servers, switches, and networking equipment co-located at the Colombo facility. Eliminates ongoing cloud costs, provides full data sovereignty, and supports 5M+ user scale. Estimated infrastructure buildout: $2.5M of the $5M raise.",
         Inches(0.9), Inches(5.88), Inches(11.5), Inches(1.1),
         font_size=13, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Technology & Infrastructure", Inches(0.85), Inches(0.3), Inches(11), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "Enterprise-grade stack — built to scale from 1,000 to 50 million users without architectural rewrites.",
         Inches(0.85), Inches(1.2), Inches(11.5), Inches(0.45),
         font_size=15, color=LIGHTGRAY, italic=True)

tech_cols = [
    ("CURRENT STACK", GOLD, [
        "Frontend: React 18, Vite, Tailwind CSS",
        "Backend: Deno / TypeScript serverless functions",
        "Database: Base44 managed (PostgreSQL)",
        "Video CDN: Cloudflare Stream",
        "Live Streaming: Cloudflare Calls SFU (WebRTC)",
        "Payments: Stripe (coins + gifting)",
        "Auth: Google OAuth + email/OTP",
        "Hosting: Vercel (web) + Base44 (API)",
        "Mobile: React Native (Android + iOS)",
    ]),
    ("SCALE INFRASTRUCTURE\n(Phase 3 — Colombo)", CORAL, [
        "Bare metal servers: Dell PowerEdge / HPE ProLiant",
        "Networking: Cisco / Juniper switching fabric",
        "Video ingest: On-prem RTMP + HLS origin servers",
        "CDN: Hybrid cloud + bare metal egress",
        "Database: Distributed PostgreSQL cluster",
        "Storage: Ceph or MinIO object storage",
        "Security: Hardware firewall, DDoS mitigation",
        "Monitoring: Prometheus + Grafana stack",
        "Compliance: GDPR, CCPA, data sovereignty",
    ]),
    ("COMPETITIVE ADVANTAGES", RGBColor(0x7B, 0xD4, 0xFF), [
        "AI moderation built-in — no third-party cost",
        "Coin economy fully proprietary — no rev share",
        "Live + podcasting + video in ONE app",
        "Location-aware content discovery",
        "Diaspora-first content strategy",
        "US-owned, US-incorporated (LDNA Consulting LLC)",
        "Trademark protected brand — USPTO filed",
        "No TikTok ban risk — fully independent",
        "Sri Lanka ops = 6–8x US cost efficiency",
    ]),
]

for i, (title, color, items) in enumerate(tech_cols):
    left = Inches(0.7 + i * 4.1)
    add_text(s, title, left, Inches(1.75), Inches(3.85), Inches(0.55),
             font_size=13, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(s, "· " + item, left, Inches(2.4 + j * 0.47), Inches(3.9), Inches(0.42),
                 font_size=12, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — USE OF FUNDS
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Use of Funds — $5M Seed Round", Inches(0.85), Inches(0.3), Inches(11), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

categories = [
    ("Infrastructure & Bare Metal",  "$2,500,000",  "50%",  GOLD,
     "Colombo bare metal servers, switches, networking\nhardware, co-location facility setup, redundancy"),
    ("Engineering Team (Colombo)",   "$1,200,000",  "24%",  CORAL,
     "10–15 developers, QA, DevOps — 24 months\nSenior engineering director + team leads"),
    ("Platform Development",         "$600,000",    "12%",  RGBColor(0x7B, 0xD4, 0xFF),
     "iOS/Android app store launch, feature roadmap\nAPI development, third-party integrations"),
    ("Marketing & User Acquisition", "$450,000",    "9%",   RGBColor(0xA0, 0xFF, 0xB0),
     "Creator fund seed, influencer partnerships\nPaid acquisition, diaspora community outreach"),
    ("Operations & Legal",           "$250,000",    "5%",   LIGHTGRAY,
     "US operations, trademark filings, legal\nCompliance, accounting, executive team"),
]

for i, (cat, amount, pct, color, desc) in enumerate(categories):
    top = Inches(1.45 + i * 0.98)
    # Bar background
    bar_bg = s.shapes.add_shape(1, Inches(0.7), top + Inches(0.28), Inches(8.5), Inches(0.35))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = MIDNAVY
    bar_bg.line.fill.background()
    # Bar fill
    bar_w = Inches(8.5 * float(pct.strip('%')) / 100)
    bar_fill = s.shapes.add_shape(1, Inches(0.7), top + Inches(0.28), bar_w, Inches(0.35))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = color
    bar_fill.line.fill.background()

    add_text(s, cat, Inches(0.7), top, Inches(5.5), Inches(0.32),
             font_size=14, bold=True, color=color)
    add_text(s, f"{amount}  ({pct})", Inches(6.3), top, Inches(3.0), Inches(0.32),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(s, desc, Inches(9.35), top + Inches(0.08), Inches(3.7), Inches(0.55),
             font_size=11, color=LIGHTGRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — FINANCIAL PROJECTIONS
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Financial Projections", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "Conservative model — based on 2% coin conversion rate and 0.5% paid subscriber rate at scale.",
         Inches(0.85), Inches(1.2), Inches(11.5), Inches(0.4),
         font_size=14, italic=True, color=LIGHTGRAY)

# Table headers
headers = ["Metric", "Year 1 (2026)", "Year 2 (2027)", "Year 3 (2028)", "Year 5 (2030)"]
rows = [
    ["MAU",                    "10,000",      "150,000",     "500,000",     "3,000,000"],
    ["Coin Purchasers (2%)",   "200",         "3,000",       "10,000",      "60,000"],
    ["Avg Spend / Purchaser",  "$25/mo",      "$30/mo",      "$35/mo",      "$40/mo"],
    ["Gifting Revenue",        "$60K",        "$1.1M",       "$4.2M",       "$28.8M"],
    ["Subscription Revenue",   "$0",          "$180K",       "$900K",       "$7.2M"],
    ["Ad Revenue",             "$0",          "$75K",        "$500K",       "$6.0M"],
    ["TOTAL REVENUE",          "$60K",        "$1.35M",      "$5.6M",       "$42M"],
    ["Gross Margin",           "—",           "55%",         "65%",         "72%"],
    ["EBITDA",                 "($3.2M)",     "($1.8M)",     "$0.8M",       "$18M+"],
]

col_widths = [Inches(2.3), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
col_lefts  = [Inches(0.7), Inches(3.1), Inches(5.1), Inches(7.1), Inches(9.1)]
row_height = Inches(0.46)
row_top_start = Inches(1.75)

for ci, header in enumerate(headers):
    box = s.shapes.add_shape(1, col_lefts[ci], row_top_start, col_widths[ci], row_height)
    box.fill.solid()
    box.fill.fore_color.rgb = GOLD
    box.line.color.rgb = NAVY
    box.line.width = Pt(0.5)
    add_text(s, header, col_lefts[ci] + Inches(0.05), row_top_start + Inches(0.07),
             col_widths[ci] - Inches(0.1), row_height,
             font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    top = row_top_start + row_height * (ri + 1)
    is_total = "TOTAL" in row[0] or "EBITDA" in row[0]
    row_color = MIDNAVY if ri % 2 == 0 else RGBColor(0x1A, 0x1C, 0x38)
    for ci, cell in enumerate(row):
        box = s.shapes.add_shape(1, col_lefts[ci], top, col_widths[ci], row_height)
        box.fill.solid()
        box.fill.fore_color.rgb = row_color
        box.line.color.rgb = RGBColor(0x30, 0x32, 0x50)
        box.line.width = Pt(0.5)
        txt_color = GOLD if is_total else (GOLD if ci == 0 else WHITE)
        add_text(s, cell, col_lefts[ci] + Inches(0.05), top + Inches(0.07),
                 col_widths[ci] - Inches(0.1), row_height,
                 font_size=12, bold=is_total, color=txt_color,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Competitive Landscape", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

competitors = [
    ("TikTok",       "✓", "✗", "✗", "✗", "✗", "US ban risk, Chinese-owned"),
    ("Instagram",    "✓", "✗", "✗", "Partial", "✗", "No coins, weak creator pay"),
    ("YouTube",      "✓", "✓", "Partial", "✗", "✗", "No short-form gifting"),
    ("Twitch",       "✗", "✓", "✓", "✗", "✗", "Gaming-only, no video feed"),
    ("Spotify",      "✗", "✗", "✓", "✗", "✗", "Audio only, no social layer"),
    ("Sachi Stream", "✓", "✓", "✓", "✓", "✓", "Full stack — only platform\nwith all 5 features"),
]

headers2 = ["Platform", "Short Video", "Live Gifting", "Podcasting", "Coin Economy", "Diaspora Focus", "Notes"]
col_w2 = [Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.7), Inches(2.55)]
col_l2 = [Inches(0.5)]
for w in col_w2[:-1]:
    col_l2.append(col_l2[-1] + w)

rh2 = Inches(0.5)
rt2 = Inches(1.3)

for ci, h in enumerate(headers2):
    box = s.shapes.add_shape(1, col_l2[ci], rt2, col_w2[ci], rh2)
    box.fill.solid()
    box.fill.fore_color.rgb = GOLD
    box.line.color.rgb = NAVY
    box.line.width = Pt(0.5)
    add_text(s, h, col_l2[ci] + Inches(0.05), rt2 + Inches(0.08),
             col_w2[ci] - Inches(0.1), rh2,
             font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

for ri, row in enumerate(competitors):
    top2 = rt2 + rh2 * (ri + 1)
    is_sachi = row[0] == "Sachi Stream"
    rc = RGBColor(0x10, 0x18, 0x10) if is_sachi else (MIDNAVY if ri % 2 == 0 else RGBColor(0x1A, 0x1C, 0x38))
    for ci, cell in enumerate(row):
        box = s.shapes.add_shape(1, col_l2[ci], top2, col_w2[ci], rh2)
        box.fill.solid()
        box.fill.fore_color.rgb = rc
        box.line.color.rgb = RGBColor(0x30, 0x32, 0x50)
        box.line.width = Pt(0.5)
        cell_color = GOLD if is_sachi else (CORAL if cell == "✗" else (RGBColor(0x7B, 0xD4, 0xFF) if cell == "✓" else WHITE))
        add_text(s, cell, col_l2[ci] + Inches(0.05), top2 + Inches(0.08),
                 col_w2[ci] - Inches(0.1), rh2,
                 font_size=12 if ci != 0 else 13, bold=is_sachi,
                 color=cell_color, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — FOUNDING TEAM
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "Founding Team", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

team_members = [
    ("Dhananjaya (Jay) Gunaratne", "Founder & CEO", GOLD,
     "21+ years Senior Infrastructure Project Manager\n"
     "Nokia, MPLS, NERC-CIP, SCADA, data center, telecom\n"
     "PSEG Long Island — $10M+ infrastructure programs\n"
     "Founder, LDNA Consulting LLC\n"
     "Architect of Sachi Stream platform vision and strategy"),
    ("Aanjinie Gunaratne", "Head of Creator & Community", CORAL,
     "Platform creator and founding community lead\n"
     "Sachi content strategy and creator culture\n"
     "User experience design and community onboarding\n"
     "Brand identity and diaspora community engagement"),
    ("Colombo Engineering Director", "VP Engineering (To Be Hired)", RGBColor(0x7B, 0xD4, 0xFF),
     "Target: Senior full-stack engineering leader\n"
     "Based in Colombo, Sri Lanka\n"
     "Will manage 10–15 developer team\n"
     "Responsible for platform migration and scaling\n"
     "Recruitment underway — Q3 2026 target start"),
]

for i, (name, title, color, bio) in enumerate(team_members):
    left = Inches(0.7 + i * 4.15)
    box = s.shapes.add_shape(1, left, Inches(1.4), Inches(3.85), Inches(4.9))
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    # Avatar circle
    circ = s.shapes.add_shape(9, left + Inches(1.45), Inches(1.55), Inches(0.95), Inches(0.95))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_text(s, name[0], left + Inches(1.55), Inches(1.7), Inches(0.75), Inches(0.65),
             font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s, name,  left + Inches(0.15), Inches(2.65), Inches(3.55), Inches(0.45),
             font_size=14, bold=True, color=color)
    add_text(s, title, left + Inches(0.15), Inches(3.1), Inches(3.55), Inches(0.35),
             font_size=12, italic=True, color=LIGHTGRAY)
    add_text(s, bio,   left + Inches(0.15), Inches(3.55), Inches(3.55), Inches(2.5),
             font_size=11, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — THE ASK
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s)
left_bar(s, GOLD)

add_text(s, "The Ask", Inches(0.85), Inches(0.3), Inches(10), Inches(0.65),
         font_size=36, bold=True, color=GOLD)
divider(s, Inches(1.1))

add_text(s, "$5,000,000", Inches(0.85), Inches(1.3), Inches(11.5), Inches(1.0),
         font_size=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Seed Round  ·  LDNA Consulting LLC  ·  May 2026",
         Inches(0.85), Inches(2.3), Inches(11.5), Inches(0.45),
         font_size=18, color=LIGHTGRAY, align=PP_ALIGN.CENTER)

divider(s, Inches(2.9), CORAL)

cols = [
    ("What We're Offering", GOLD, [
        "Equity stake — percentage TBD in negotiation",
        "Board observer seat available at $1M+ commitment",
        "Series A right of first refusal",
        "Quarterly financial reporting",
        "Direct access to founding team",
    ]),
    ("What We'll Deliver in 18 Months", CORAL, [
        "iOS + Android app live on both stores",
        "Colombo engineering team operational",
        "150,000+ MAU target",
        "Bare metal infrastructure Phase 1 online",
        "$1M+ ARR run rate",
    ]),
    ("Path to Series A", RGBColor(0x7B, 0xD4, 0xFF), [
        "Series A target: $10–15M at 250K+ MAU",
        "Expected timing: 18–24 months post-seed",
        "Target valuation: $40–60M pre-money",
        "Lead investor for Series A being identified",
        "Secondary markets: South Asia, Pacific",
    ]),
]

for i, (title, color, items) in enumerate(cols):
    left = Inches(0.7 + i * 4.15)
    box = s.shapes.add_shape(1, left, Inches(3.2), Inches(3.85), Inches(3.55))
    box.fill.solid()
    box.fill.fore_color.rgb = MIDNAVY
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_text(s, title, left + Inches(0.15), Inches(3.3), Inches(3.55), Inches(0.4),
             font_size=13, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(s, "· " + item, left + Inches(0.15), Inches(3.85 + j * 0.5),
                 Inches(3.55), Inches(0.45), font_size=12, color=SOFTWHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSE
# ══════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
accent_bar(s, GOLD, height=Inches(0.05), top=Inches(0))
accent_bar(s, GOLD, height=Inches(0.05), top=Inches(7.45))

add_text(s, "Let's Build the Future\nof Creator Culture Together.",
         Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.0),
         font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

divider(s, Inches(3.7), CORAL)

add_text(s, "Dhananjaya (Jay) Gunaratne\nFounder & CEO, LDNA Consulting LLC",
         Inches(0.8), Inches(3.95), Inches(11.5), Inches(0.75),
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "jaygnz27@gmail.com  ·  908-255-2195  ·  sachistream.com",
         Inches(0.8), Inches(4.75), Inches(11.5), Inches(0.45),
         font_size=16, color=LIGHTGRAY, align=PP_ALIGN.CENTER)

add_text(s, "New Providence, NJ 07974  ·  LDNA Consulting LLC",
         Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
         font_size=13, color=LIGHTGRAY, align=PP_ALIGN.CENTER, italic=True)

divider(s, Inches(5.8))

add_text(s, "\"Sachi — happiness, good fortune, and bliss.\nWe're building the platform where creators thrive.\"",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.75),
         font_size=15, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
filename = "Sachi_Stream_Investor_Deck_Seed_Round_May2026.pptx"
prs.save(filename)
print(f"Saved: {filename}")
